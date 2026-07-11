"""
Text extraction for uploaded RAG documents.

The extractor supports plain text/code formats without extra dependencies and
uses optional parser libraries for richer formats when installed.
"""

from __future__ import annotations

import csv
import json
from pathlib import Path

from bs4 import BeautifulSoup

TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".rst",
    ".log",
    ".py",
    ".js",
    ".ts",
    ".jsx",
    ".tsx",
    ".html",
    ".htm",
    ".css",
    ".json",
    ".jsonl",
    ".yml",
    ".yaml",
    ".toml",
    ".ini",
    ".cfg",
    ".conf",
    ".xml",
    ".sql",
    ".graphql",
    ".proto",
    ".sh",
    ".bash",
    ".zsh",
    ".c",
    ".cpp",
    ".h",
    ".hpp",
    ".java",
    ".go",
    ".rs",
    ".rb",
    ".php",
}
TABLE_EXTENSIONS = {".csv", ".tsv"}
SUPPORTED_EXTENSIONS = (
    TEXT_EXTENSIONS
    | TABLE_EXTENSIONS
    | {
        ".pdf",
        ".docx",
        ".xlsx",
        ".xlsm",
        ".pptx",
        ".ipynb",
    }
)


class UnsupportedFileTypeError(ValueError):
    """Raised when a file extension cannot be indexed."""


class TextExtractionError(RuntimeError):
    """Raised when a supported file cannot be converted to text."""


def supported_extensions() -> list[str]:
    """Return supported extensions for API responses and frontend hints."""
    return sorted(SUPPORTED_EXTENSIONS)


def extract_text(path: str | Path, original_filename: str | None = None) -> str:
    """Extract indexable text from an uploaded file."""
    file_path = Path(path)
    ext_source = original_filename or file_path.name
    ext = Path(ext_source).suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        raise UnsupportedFileTypeError(f"Unsupported file type '{ext or 'unknown'}'.")

    if ext in TEXT_EXTENSIONS:
        text = file_path.read_text(encoding="utf-8", errors="ignore")
        if ext in {".html", ".htm"}:
            return BeautifulSoup(text, "html.parser").get_text("\n")
        if ext == ".ipynb":
            return _extract_ipynb_text(text)
        return text

    if ext in TABLE_EXTENSIONS:
        return _extract_delimited_text(
            file_path, delimiter="\t" if ext == ".tsv" else ","
        )

    if ext == ".pdf":
        return _extract_pdf_text(file_path)

    if ext == ".docx":
        return _extract_docx_text(file_path)

    if ext in {".xlsx", ".xlsm"}:
        return _extract_xlsx_text(file_path)

    if ext == ".pptx":
        return _extract_pptx_text(file_path)

    raise UnsupportedFileTypeError(f"Unsupported file type '{ext or 'unknown'}'.")


def _extract_ipynb_text(raw: str) -> str:
    try:
        notebook = json.loads(raw)
    except json.JSONDecodeError as exc:
        raise TextExtractionError(f"Invalid notebook JSON: {exc}") from exc

    parts = []
    for cell in notebook.get("cells", []):
        source = cell.get("source", "")
        if isinstance(source, list):
            source = "".join(source)
        if source:
            parts.append(str(source))
    return "\n\n".join(parts)


def _extract_delimited_text(path: Path, delimiter: str) -> str:
    rows = []
    with path.open("r", encoding="utf-8", errors="ignore", newline="") as f:
        reader = csv.reader(f, delimiter=delimiter)
        for row in reader:
            rows.append(" | ".join(cell.strip() for cell in row))
    return "\n".join(rows)


def _extract_pdf_text(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise TextExtractionError(
            "PDF upload support requires the 'pypdf' package."
        ) from exc

    reader = PdfReader(str(path))
    pages = []
    for page in reader.pages:
        pages.append(page.extract_text() or "")
    return "\n\n".join(pages)


def _extract_docx_text(path: Path) -> str:
    try:
        from docx import Document as DocxDocument
    except ImportError as exc:
        raise TextExtractionError(
            "DOCX upload support requires the 'python-docx' package."
        ) from exc

    doc = DocxDocument(str(path))
    parts = [p.text for p in doc.paragraphs if p.text]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(parts)


def _extract_xlsx_text(path: Path) -> str:
    try:
        import openpyxl
    except ImportError as exc:
        raise TextExtractionError(
            "XLSX upload support requires the 'openpyxl' package."
        ) from exc

    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet in workbook.worksheets:
        parts.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                parts.append(" | ".join(values))
    workbook.close()
    return "\n".join(parts)


def _extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise TextExtractionError(
            "PPTX upload support requires the 'python-pptx' package."
        ) from exc

    presentation = Presentation(str(path))
    parts = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"# Slide {slide_index}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)
